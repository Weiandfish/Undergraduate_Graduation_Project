"""Build a polished graduation defense PPT using BIT template layouts."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import io, zipfile, os, copy

# ============================================================
# STEP 1: Clean template - remove all old slides
# ============================================================
TEMPLATE = '北京理工大学学术答辩PPT模板V3.0.pptx'
OUTPUT = '毕业答辩_唐元昊_v2.pptx'

print("Cleaning template...")
with open(TEMPLATE, 'rb') as f:
    template_bytes = f.read()

zip_in = zipfile.ZipFile(io.BytesIO(template_bytes), 'r')
zip_out = io.BytesIO()

with zipfile.ZipFile(zip_out, 'w', zipfile.ZIP_DEFLATED) as zout:
    for item in zip_in.infolist():
        if item.filename.startswith('ppt/slides/slide') and 'slideM' not in item.filename:
            continue
        if item.filename.startswith('ppt/slides/_rels/slide') and 'slideM' not in item.filename:
            continue
        if item.filename == 'ppt/presentation.xml':
            data = zip_in.read(item)
            root = etree.fromstring(data)
            nsmap = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            }
            sldIdLst = root.find('.//p:sldIdLst', nsmap)
            if sldIdLst is not None:
                for child in list(sldIdLst):
                    sldIdLst.remove(child)
            data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
            zout.writestr(item, data)
        elif item.filename == 'ppt/_rels/presentation.xml.rels':
            data = zip_in.read(item)
            root = etree.fromstring(data)
            for rel in list(root):
                target = rel.get('Target', '')
                if target.startswith('slides/slide') and 'slideM' not in target:
                    root.remove(rel)
            data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
            zout.writestr(item, data)
        else:
            zout.writestr(item, zip_in.read(item))
zip_in.close()

prs = Presentation(io.BytesIO(zip_out.getvalue()))
print(f"Cleaned. Slides: {len(prs.slides)}")

# ============================================================
# STEP 2: Layout references
# ============================================================
# Master 0: Style 1 title/end pages
LAYOUT_TITLE = prs.slide_masters[0].slide_layouts[0]   # 标题式1-首页
LAYOUT_END   = prs.slide_masters[0].slide_layouts[1]   # 标题式1-尾页

# Master 11: 目录样式3-2 - large number section divider
LAYOUT_SECTION = prs.slide_masters[11].slide_layouts[1]  # 目录样式3-2

# Master 17: Content pages
M17 = prs.slide_masters[17]
LAYOUT_TEXT   = M17.slide_layouts[0]   # 内容页样式1-纯文
LAYOUT_IMG1   = M17.slide_layouts[1]   # 内容页样式1-一图一图-1 (image top, text below)
LAYOUT_IMG2   = M17.slide_layouts[2]   # 内容页样式1-一图一图-2 (image left, text right)

print("Layouts ready.")

# ============================================================
# Colors
# ============================================================
GREEN    = RGBColor(0, 108, 57)
DARK     = RGBColor(63, 63, 63)
GRAY     = RGBColor(162, 162, 162)
RED      = RGBColor(161, 63, 11)
WHITE    = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(245, 248, 246)   # Very light green tint
BORDER   = RGBColor(220, 230, 224)   # Subtle border

# ============================================================
# Helper functions
# ============================================================
def set_title(slide, text):
    """Set the title placeholder text with proper formatting."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = '微软雅黑'
                    run.font.size = Pt(28)
                    run.font.bold = True
                    try:
                        run.font.color.rgb = DARK
                    except:
                        pass
            break

def add_body(slide, left, top, width, height, items):
    """Add body text box with proper spacing and typography."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if isinstance(item, tuple):
            text, fs, b, c = item
            p.text = text
            p.font.size = fs
            p.font.bold = b
            p.font.color.rgb = c
        else:
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = DARK

        p.font.name = '微软雅黑'

        # Line spacing
        pPr = p._p.get_or_add_pPr()
        # Space before/after
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPtsBef = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPtsBef.set('val', str(200 if i > 0 and isinstance(item, tuple) and item[2] else 0))

        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spcPtsAft = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPtsAft.set('val', '100')

        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', '160000')

    return txBox

def add_subtitle(slide, text, left=1606550, top=950000):
    """Add a subtitle/English text below the main title."""
    return add_body(slide, left, top, 8000000, 400000, [
        (text, Pt(12), False, GRAY)
    ])

def add_image_placeholder(slide, left, top, width, height, label="[ 图片占位 ]"):
    """Add a visible placeholder rectangle for images."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    # Center text vertically
    txBody = shape._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')

    return shape

def add_section_divider(slide, num, title_cn, title_en):
    """Build a section divider slide using 目录样式3-2 layout."""
    # This layout has shapes for large number, subtitle, line
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            full_text = tf.text
            if full_text.strip() == '05' or full_text.strip().startswith('0'):
                # Set section number
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.text = ''
                tf.paragraphs[0].text = num
                for run in tf.paragraphs[0].runs:
                    run.font.size = Pt(120)
                    run.font.bold = True
                    try:
                        run.font.color.rgb = GREEN
                    except:
                        pass
            elif '题' in full_text or 'Background' in full_text or '项目' in full_text:
                # Set title
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.text = ''
                tf.paragraphs[0].text = title_cn
                for run in tf.paragraphs[0].runs:
                    run.font.name = '微软雅黑'
                    run.font.size = Pt(32)
                    run.font.bold = True
                    try:
                        run.font.color.rgb = DARK
                    except:
                        pass
                # Add English subtitle
                p2 = tf.add_paragraph()
                p2.text = title_en
                for run in p2.runs:
                    run.font.name = 'Arial'
                    run.font.size = Pt(14)
                    try:
                        run.font.color.rgb = GRAY
                    except:
                        pass

    # Also try to add title as a new textbox if shapes were not found
    found_text = False
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            found_text = True
            break

    if not found_text:
        # Add text manually as fallback
        add_body(slide, 838200, 2000000, 9000000, 1200000, [
            (num, Pt(96), True, GREEN),
            ('', Pt(8), False, DARK),
            (title_cn, Pt(32), True, DARK),
            (title_en, Pt(14), False, GRAY),
        ])

    # Add page number at bottom right
    add_body(slide, 10000000, 6200000, 2000000, 300000, [
        ('', Pt(10), False, GRAY)
    ])


# ============================================================
# BUILD SLIDES
# ============================================================
TOTAL = 17

# === Slide 1: Title Page ===
print("Building...")
slide = prs.slides.add_slide(LAYOUT_TITLE)
for shape in slide.placeholders:
    ph = shape.placeholder_format
    if ph.idx == 12:
        shape.text = '面向桌面操作的VLA机械臂\n抓取与控制研究'
        for p in shape.text_frame.paragraphs:
            p.font.size = Pt(40)
            p.font.bold = True
            try: p.font.color.rgb = WHITE
            except: pass
            p.alignment = PP_ALIGN.CENTER
    elif ph.idx == 13:
        shape.text = '答辩人：唐元昊      指导教师：方浩 教授\n北京理工大学 自动化学院'
        for p in shape.text_frame.paragraphs:
            p.font.size = Pt(16)
            try: p.font.color.rgb = WHITE
            except: pass
            p.alignment = PP_ALIGN.CENTER
print("  Slide 1: Title")

# === Slide 2: TOC ===
slide = prs.slides.add_slide(LAYOUT_TEXT)
set_title(slide, '目  录  |  CONTENTS')

toc_items = [
    ('01', '研究背景与文献综述'),
    ('02', 'Franka系统集成与LeRobot接口扩展'),
    ('03', '桌面抓取示教数据采集与数据集规范'),
    ('04', 'π₀.₅模型训练与实验分析'),
    ('05', '总结与展望'),
]
y = 1550000
for num, title in toc_items:
    add_body(slide, 2500000, y, 8000000, 600000, [
        (f'{num}    {title}', Pt(22), False, DARK),
    ])
    # Green number accent on left
    add_body(slide, 1200000, y, 900000, 600000, [
        (num, Pt(28), True, GREEN),
    ])
    y += 700000
    # subtle separator
    if num != '05':
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(2500000), Emu(y - 60000), Emu(8000000), Emu(8000))
        sep.fill.solid()
        sep.fill.fore_color.rgb = BORDER
        sep.line.fill.background()
add_body(slide, 10000000, 6350000, 1800000, 300000, [
    (f'2 / {TOTAL}', Pt(9), False, GRAY)
])
print("  Slide 2: TOC")

# === Slide 3: Section 1 Divider ===
slide = prs.slides.add_slide(LAYOUT_SECTION)
add_section_divider(slide, '01', '研究背景与文献综述', 'Background & Literature Review')
# Add decorative green line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(838200), Emu(2550000), Emu(2500000), Emu(40000))
line.fill.solid(); line.fill.fore_color.rgb = GREEN; line.line.fill.background()
add_body(slide, 10000000, 6200000, 2000000, 300000, [(f'3 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 3: Section 1 divider")

# === Slide 4: Research Background ===
slide = prs.slides.add_slide(LAYOUT_TEXT)
set_title(slide, '研究背景')
add_subtitle(slide, 'Research Background')
add_body(slide, 1606550, 1100000, 9400000, 5000000, [
    ('VLA模型的核心突破', Pt(18), True, GREEN),
    '视觉-语言-动作（VLA）模型将多模态大模型的语义理解与机器人动作生成相统一',
    '为"自然语言条件下的灵巧操作"提供了全新的技术路径',
    '',
    ('工程落地的现实挑战', Pt(18), True, GREEN),
    '从开源算法（π₀.₅、OpenVLA）到具体实验室硬件（Franka机械臂、多路相机、实时控制链路）',
    '中间普遍存在接口割裂、时序对齐困难、示教可重复性差等问题',
    '桌面抓取任务同时考验全局场景理解与接触阶段的精细反馈 — 仅靠"更换超参数"无法解决',
    '',
    ('核心命题', Pt(18), True, GREEN),
    '如何建立从硬件接入、数据采集、模型训练到实机部署的可复现全栈闭环？',
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'4 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 4: Research background")

# === Slide 5: Literature Review ===
slide = prs.slides.add_slide(LAYOUT_TEXT)
set_title(slide, '文献综述：VLA技术演进路线')
add_subtitle(slide, 'Literature Review — Evolution of VLA Models')
add_body(slide, 1606550, 1100000, 9400000, 5000000, [
    ('早期探索 (2021-2022)', Pt(16), True, GREEN),
    'CLIPort 语言-视觉对齐抓取 → SayCan LLM任务规划+价值筛选 → VIMA 多模态提示',
    '',
    ('具身基础模型 (2023)', Pt(16), True, GREEN),
    'PaLM-E 感知流注入LLM实现长程规划 → RT-1/RT-2 大规模真机数据训练',
    'Open X-Embodiment 联合多实验室训练RT-X — 跨硬件统计迁移成为共识',
    '',
    ('开源VLA时代 (2024-2025)', Pt(16), True, GREEN),
    'OpenVLA / Octo — 可本地微调、可替换相机与语言条件的通用策略',
    'π₀ / π₀.₅ — 分层推理架构：高层语义子任务 + Flow Matching动作专家',
    'LeRobot — 统一张量规范，推动跨实验室数据共享与复现',
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'5 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 5: Literature review")

# === Slide 6: Research Hypotheses ===
slide = prs.slides.add_slide(LAYOUT_TEXT)
set_title(slide, '研究问题与研究假设')
add_subtitle(slide, 'Research Questions & Hypotheses')
add_body(slide, 1606550, 1100000, 9400000, 5200000, [
    ('H1  跨平台迁移可行性', Pt(20), True, GREEN),
    'π₀.₅能否从原论文的移动双臂设定迁移至静态Franka单臂桌面抓取？',
    '适配后策略是否能学习有效的"接近 → 夹取 → 抬升 → 放置"多阶段行为？',
    '',
    ('H2  闭环策略的决定性作用', Pt(20), True, GREEN),
    '在线闭环执行策略是否是接触敏感任务的决定性因素？',
    '异步开环与分段同步闭环在同一模型权重下的成功率是否有数量级差异？',
    '',
    ('H3  遥操作标定的可重复性', Pt(20), True, GREEN),
    '同构主从遥操作标定能否实现可重复、可版本化的示教数据采集？',
    '两阶段配对标定流程能否自动解算关节缩放/偏置与夹爪量程？',
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'6 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 6: Hypotheses")

# === Slide 7: Section 2 Divider ===
slide = prs.slides.add_slide(LAYOUT_SECTION)
add_section_divider(slide, '02', 'Franka系统集成与\nLeRobot接口扩展', 'System Integration & LeRobot Interface')
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(838200), Emu(2550000), Emu(2500000), Emu(40000))
line.fill.solid(); line.fill.fore_color.rgb = GREEN; line.line.fill.background()
add_body(slide, 10000000, 6200000, 2000000, 300000, [(f'7 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 7: Section 2 divider")

# === Slide 8: System Architecture (with image placeholder) ===
slide = prs.slides.add_slide(LAYOUT_IMG1)
set_title(slide, '系统总体架构 — 四层设计')
add_subtitle(slide, 'Four-Layer System Architecture')
# Image placeholder at top (the layout has image area at ~0-3000000)
add_image_placeholder(slide, 500000, 1100000, 11000000, 2800000,
                      '[ 系统架构图占位 — 四层架构示意图 ]\n\n实时执行层 | 网络服务层 | LeRobot抽象层 | 遥操作数据层')
# Text description below
add_body(slide, 1606550, 4100000, 9000000, 2000000, [
    ('四层架构：实时控制(C++ franka_server) → TCP套接字通信 → LeRobot Robot抽象 → 遥操作接口', Pt(15), False, DARK),
    ('核心创新：主从映射JSON参数化 + 两阶段交互式配对标定，标定结果随数据集版本化', Pt(15), False, GREEN),
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'8 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 8: Architecture")

# === Slide 9: Hardware & Teleoperation (image left + text right) ===
slide = prs.slides.add_slide(LAYOUT_IMG2)
set_title(slide, '硬件平台与遥操作方案')
add_subtitle(slide, 'Hardware Platform & Teleoperation')

# Clean up layout text placeholders and add our content
# Layout 2 has image area on left side
add_image_placeholder(slide, 400000, 1200000, 5600000, 3800000,
                      '[ 硬件平台照片占位 ]\n\nFranka FR3 + 夹爪 + 相机布置\n同构主从臂示教场景')
add_body(slide, 6300000, 1200000, 5200000, 5000000, [
    ('硬件配置', Pt(18), True, GREEN),
    'Franka FR3 七轴机械臂 (±0.1mm)',
    '知行 CTAG2F90C 二指平行夹爪',
    'Intel RealSense D445 腕部深度相机',
    'RTX A6000 (48GB) + Xeon Platinum',
    '',
    ('遥操作方案', Pt(18), True, GREEN),
    'SpaceMouse：6-DOF 笛卡尔增量 → IK',
    '同构主从：theta_f = scale·theta_l + bias',
    '两阶段标定：关节对齐 → 夹爪量程',
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'9 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 9: Hardware")

# === Slide 10: Section 3 Divider ===
slide = prs.slides.add_slide(LAYOUT_SECTION)
add_section_divider(slide, '03', '桌面抓取示教数据采集\n与LeRobot数据集规范', 'Data Collection & LeRobot Standard')
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(838200), Emu(2550000), Emu(2500000), Emu(40000))
line.fill.solid(); line.fill.fore_color.rgb = GREEN; line.line.fill.background()
add_body(slide, 10000000, 6200000, 2000000, 300000, [(f'10 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 10: Section 3 divider")

# === Slide 11: Dataset (text + image) ===
slide = prs.slides.add_slide(LAYOUT_IMG1)
set_title(slide, 'pick_apple_4_18 数据集')
add_subtitle(slide, 'LeRobot v2.1 Compliant — Desktop Apple Grasping Dataset')

# Image placeholder: visualization screenshots
add_image_placeholder(slide, 600000, 1100000, 5400000, 2400000,
                      '[ Rerun可视化截图 ]\n关节位置时序 · 末端轨迹 · 视频同步')
add_image_placeholder(slide, 6200000, 1100000, 5400000, 2400000,
                      '[ 数据采集场景 ]\n双相机视角：第三人称 + 腕部D445')

# Key specs
add_body(slide, 1606550, 3700000, 9000000, 2800000, [
    ('任务：Pick up the green apple and place it in the upper left corner.', Pt(15), True, DARK),
    ('30 episodes | 5,512 帧 | 60条MP4视频 | 10Hz采样 | 8维状态/动作 | 双路RGB 1280×720', Pt(14), False, DARK),
    '',
    ('数据格式：meta/ (JSON+JSONL元数据)  +  data/ (Parquet时序表)  +  videos/ (MP4外挂视频)', Pt(14), False, GRAY),
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'11 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 11: Dataset")

# === Slide 12: Section 4 Divider ===
slide = prs.slides.add_slide(LAYOUT_SECTION)
add_section_divider(slide, '04', '基于π₀.₅的抓取策略\n训练与实验分析', 'pi0.5 Training & Experimental Analysis')
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(838200), Emu(2550000), Emu(2500000), Emu(40000))
line.fill.solid(); line.fill.fore_color.rgb = GREEN; line.line.fill.background()
add_body(slide, 10000000, 6200000, 2000000, 300000, [(f'12 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 12: Section 4 divider")

# === Slide 13: pi0.5 Architecture ===
slide = prs.slides.add_slide(LAYOUT_IMG2)
set_title(slide, 'π₀.₅模型核心原理')
add_subtitle(slide, 'Hierarchical VLA with Flow Matching Action Expert')

# Left: model architecture placeholder
add_image_placeholder(slide, 400000, 1200000, 5600000, 4000000,
                      '[ π₀.₅模型架构图占位 ]\n\nVLM主干：SigLIP + Gemma\n动作专家：Flow Matching\nadaRMSNorm时间注入')

# Right: key points
add_body(slide, 6300000, 1200000, 5200000, 5000000, [
    ('分层推理', Pt(18), True, GREEN),
    'π(a, l̂|o, l) = π(l̂|o, l) · π(a|o, l̂)',
    '高层：推断语义子任务（"抓住苹果"→"放到左上角"）',
    '低层：Flow Matching 生成连续平滑动作块序列',
    '',
    ('VLM主干 + Action Expert', Pt(18), True, GREEN),
    'SigLIP视觉编码器 + Gemma语言骨干',
    '双视角图像在同一Transformer中建立跨模态关联',
    'pi05=True 启用adaRMSNorm时间注入',
    '',
    ('迁移适配', Pt(18), True, GREEN),
    '移动双臂 → 静态Franka单臂',
    'third_person → images/base, wrist → images/wrist',
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'13 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 13: pi0.5 model")

# === Slide 14: Training Configuration ===
slide = prs.slides.add_slide(LAYOUT_TEXT)
set_title(slide, '训练配置与模型变体对比')
add_subtitle(slide, 'Training Configuration & Model Variants')
add_body(slide, 1606550, 1100000, 9400000, 5000000, [
    ('训练参数设置', Pt(18), True, GREEN),
    'action_horizon = 32    |    max_token_len = 200    |    batch_size = 4    |    50,000 steps',
    '30 episodes全部用于训练    |    NVIDIA RTX A6000 (48GB) ×1    |    PyTorch 2.1 + Ubuntu 22.04',
    '',
    ('三种模型变体对比', Pt(18), True, GREEN),
    '',
    ('全量微调 (Full Fine-tuning)                    成功率 ≈ 80%      显存 ≈ 12 GB', Pt(16), True, DARK),
    ('冻结VLM主干 (Frozen Backbone)            成功率 ≈ 90%      显存 ≈ 10 GB', Pt(16), True, DARK),
    ('Short Memory (减少条件帧数)               成功率 ≈ 100%    夹爪滑落后自主重新抓取', Pt(16), True, GREEN),
    '',
    ('减少条件帧数使策略更依赖即时视觉反馈，增强对接触事件的反应能力', Pt(14), False, GRAY),
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'14 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 14: Training")

# === Slide 15: KEY RESULT - Two-column comparison ===
slide = prs.slides.add_slide(LAYOUT_IMG1)
set_title(slide, '关键发现：闭环执行策略决定成败')
add_subtitle(slide, 'Key Finding — Deployment Strategy Dominates Contact-Rich Tasks')

# Two side-by-side comparison boxes
col_w = 5400000
left_x = 600000
right_x = 6200000
y_top = 1100000
box_h = 1800000

# Left: Async failure
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left_x), Emu(y_top), Emu(col_w), Emu(box_h))
left_box.fill.solid(); left_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
left_box.line.color.rgb = RGBColor(230, 200, 200); left_box.line.width = Pt(1)

add_body(slide, left_x + 200000, y_top + 100000, col_w - 400000, box_h - 200000, [
    ('异步远程客户端  ≈ 5%', Pt(22), True, RED),
    ('推理与控制线程解耦，连续执行全部32步', Pt(13), False, DARK),
    ('夹爪关闭时机偏差 → 苹果滑落 → 累积误差不可恢复', Pt(13), False, GRAY),
])

# Right: Sync success
right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(right_x), Emu(y_top), Emu(col_w), Emu(box_h))
right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(240, 250, 243)
right_box.line.color.rgb = RGBColor(200, 230, 210); right_box.line.width = Pt(1)

add_body(slide, right_x + 200000, y_top + 100000, col_w - 400000, box_h - 200000, [
    ('分段同步闭环  ≈ 100%', Pt(22), True, GREEN),
    ('仅执行第3-10步(~0.27s) → 重新观测+推理', Pt(13), False, DARK),
    ('动作块 → 滚动时域控制器，高频重规划消解累积误差', Pt(13), False, GRAY),
])

# Bottom: comparison diagram placeholder
add_image_placeholder(slide, 600000, 3100000, 11000000, 1600000,
                      '[ 对比实验示意图占位 ]\n\n异步执行：臂接近苹果 → 夹爪关闭过晚 → 苹果滑落 ✗\n同步闭环：臂接近苹果 → 短窗口执行 → 重新观测 → 动态纠偏 → 成功抓取 ✓')

# Key takeaway
add_body(slide, 1606550, 4900000, 9000000, 1200000, [
    ('核心结论', Pt(18), True, GREEN),
    ('同一模型权重，仅改变部署策略 → 成功率从 ~5% 跃升至 ~100%', Pt(16), True, DARK),
    ('闭环执行策略对接触敏感型抓取任务具有决定性作用 — 比模型架构更重要', Pt(15), False, GRAY),
])
add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'15 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 15: Key result")

# === Slide 16: Conclusion ===
slide = prs.slides.add_slide(LAYOUT_TEXT)
set_title(slide, '总结与展望')
add_subtitle(slide, 'Summary & Future Work')

# Two columns for contributions + limitations
add_body(slide, 1606550, 1100000, 5200000, 5000000, [
    ('主要贡献', Pt(20), True, GREEN),
    ('', Pt(6), False, DARK),
    ('1. Franka-LeRobot 全栈打通', Pt(16), True, DARK),
    ('   实时控制服务 + Robot/Teleoperator抽象，开源可复用', Pt(13), False, GRAY),
    ('2. 两阶段遥操作标定', Pt(16), True, DARK),
    ('   自动解算关节偏置与夹爪量程，JSON配置可版本化', Pt(13), False, GRAY),
    ('3. 标准化数据集', Pt(16), True, DARK),
    ('   pick_apple_4_18 (30ep/5512帧) LeRobot v2.1规范', Pt(13), False, GRAY),
    ('4. π₀.₅成功跨平台迁移', Pt(16), True, DARK),
    ('   移动双臂 → 静态Franka单臂，全栈适配', Pt(13), False, GRAY),
    ('5. 闭环策略关键发现', Pt(16), True, DARK),
    ('   分段同步闭环将成功率从~5%提升至~100%', Pt(13), False, GRAY),
])

add_body(slide, 7200000, 1100000, 4500000, 5000000, [
    ('局限与展望', Pt(20), True, RED),
    ('', Pt(6), False, DARK),
    ('任务分布较窄', Pt(16), True, DARK),
    ('当前仅采集单一苹果抓取数据', Pt(13), False, GRAY),
    ('缺乏跨物体、跨场景泛化', Pt(13), False, GRAY),
    ('', Pt(6), False, DARK),
    ('未来方向：', Pt(16), True, GREEN),
    ('扩展多物体/多姿态数据集', Pt(13), False, DARK),
    ('Sim-to-Real迁移 + 数据增强', Pt(13), False, DARK),
    ('更多接触富集任务的闭环策略验证', Pt(13), False, DARK),
    ('探索π₀.₅在更复杂长程任务上的表现', Pt(13), False, DARK),
])

add_body(slide, 10000000, 6350000, 1800000, 300000, [(f'16 / {TOTAL}', Pt(9), False, GRAY)])
print("  Slide 16: Conclusion")

# === Slide 17: Thanks ===
slide = prs.slides.add_slide(LAYOUT_END)
for shape in slide.placeholders:
    ph = shape.placeholder_format
    if ph.idx == 12:
        shape.text = '感谢各位老师批评指正'
        for p in shape.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            try: p.font.color.rgb = WHITE
            except: pass
            p.alignment = PP_ALIGN.CENTER
    elif ph.idx == 13:
        shape.text = '答辩人：唐元昊      指导教师：方浩 教授\n北京理工大学 自动化学院'
        for p in shape.text_frame.paragraphs:
            p.font.size = Pt(16)
            try: p.font.color.rgb = WHITE
            except: pass
            p.alignment = PP_ALIGN.CENTER
print("  Slide 17: Thanks")

# ============================================================
# SAVE
# ============================================================
print(f"\nTotal slides: {len(prs.slides)}")
prs.save(OUTPUT)
print(f"Saved to: {OUTPUT}")
print("Done!")
